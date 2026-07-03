#!/usr/bin/env python3
"""Generate the Personal Portfolio sharing PPT in minimalist white style (Scheme B)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors from spec_lock ───────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
FAFAFA      = RGBColor(0xFA, 0xFA, 0xFA)
EEEE        = RGBColor(0xEE, 0xEE, 0xEE)
F0F0F0      = RGBColor(0xF0, 0xF0, 0xF0)
CHARCOAL    = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_DARK   = RGBColor(0x66, 0x66, 0x66)
GRAY_MED    = RGBColor(0x99, 0x99, 0x99)
GRAY_LIGHT  = RGBColor(0xBB, 0xBB, 0xBB)
PURPLE      = RGBColor(0x6C, 0x63, 0xFF)
TEAL        = RGBColor(0x00, 0xD4, 0xAA)
PINK        = RGBColor(0xFF, 0x6B, 0x9D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def bg(slide, color=WHITE):
    """Set slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None, corner_radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    return shape

def add_line(slide, x1, y1, x2, y2, color=EEEE, width=Pt(1)):
    """Add a straight line."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # msoConnectorStraight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=CHARCOAL, bold=False, font_name="Microsoft YaHei",
                 alignment=PP_ALIGN.LEFT, line_spacing=1.4):
    """Add a text box with Chinese-friendly font."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_accent_circle(slide, left, top, size, color):
    """Add a small colored circle as accent dot."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ━━━━━━━━━━━━━━━━━━━━━━━ S01 Cover ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg(slide)

# Decorative circles (left side)
for cx, cy, r, alpha_color in [
    (Inches(2.5), Inches(1.8), Inches(1.2), PURPLE),
    (Inches(3.2), Inches(2.5), Inches(0.8), TEAL),
    (Inches(2.0), Inches(3.0), Inches(0.6), PINK),
]:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
    shape.fill.solid()
    shape.fill.fore_color.rgb = alpha_color
    shape.line.fill.background()
    # Low opacity via transparency hack: use semi-transparent fill
    # pptx doesn't support alpha directly, so we use a light version
    shape.fill.fore_color.rgb = alpha_color

# Title area (right side)
add_text_box(slide, Inches(5.5), Inches(2.0), Inches(6.5), Inches(1.5),
             "个人作品集", 56, CHARCOAL, True, "Microsoft YaHei", PP_ALIGN.LEFT)
add_text_box(slide, Inches(5.5), Inches(3.0), Inches(6.5), Inches(0.8),
             "我的技术探索与创造之旅", 26, GRAY_DARK, False, "Microsoft YaHei", PP_ALIGN.LEFT)
add_text_box(slide, Inches(5.5), Inches(3.7), Inches(6.5), Inches(0.5),
             "Full-Stack Developer  ·  Problem Solver  ·  Creative Coder", 16, GRAY_MED, False, "Segoe UI", PP_ALIGN.LEFT)

# Purple accent line
add_line(slide, Inches(5.5), Inches(4.4), Inches(6.2), Inches(4.4), PURPLE, Pt(3))

# Bottom text
add_text_box(slide, Inches(5.5), Inches(6.2), Inches(5), Inches(0.4),
             "2026  ·  社团技术分享", 14, GRAY_LIGHT, False, "Segoe UI", PP_ALIGN.LEFT)


# ━━━━━━━━━━━━━━━━━━━━━━━ S02 About Me ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

# Title
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "关于我 — 技术探索之旅", 32, CHARCOAL, True, "Microsoft YaHei")
add_line(slide, Inches(0.8), Inches(1.1), Inches(1.8), Inches(1.1), PURPLE, Pt(3))

# Gray container
add_rect(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(3.5), FAFAFA, EEEE, Pt(1))

add_text_box(slide, Inches(1.3), Inches(1.9), Inches(10.5), Inches(2.8),
             "我是一名热爱技术与创造的开发者。喜欢用代码构建有趣的东西，从前端界面到后端架构，\n"
             "从数据分析到人工智能，不断探索技术的边界。相信好的设计加上优雅的代码，\n"
             "就能创造出打动人心的产品。",
             20, GRAY_DARK, False, "Microsoft YaHei", PP_ALIGN.LEFT, 1.8)

# Tags
for i, (label, color) in enumerate([
    ("前端开发", PURPLE), ("后端服务", TEAL), ("数据与AI", PINK), ("结构设计", PURPLE)
]):
    x = Inches(1.3 + i * 2.9)
    shape = add_rect(slide, x, Inches(5.0), Inches(2.5), Inches(0.45), None, color, Pt(2))
    shape.text_frame.paragraphs[0].text = label
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].font.color.rgb = color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    shape.text_frame.word_wrap = True
    shape.vertical_anchor = MSO_ANCHOR.MIDDLE


# ━━━━━━━━━━━━━━━━━━━━━━━ S03 Skills Overview ━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "技术栈全景", 32, CHARCOAL, True, "Microsoft YaHei")
add_line(slide, Inches(0.8), Inches(1.1), Inches(1.8), Inches(1.1), PURPLE, Pt(3))

skills = [
    ("前端开发", "HTML / CSS / JS / React / Vue", 85, PURPLE),
    ("后端服务", "Python / Java / Node.js / REST", 75, TEAL),
    ("数据与AI", "Machine Learning / NLP / Analysis", 70, PINK),
    ("数据结构", "Algorithms / Competitive Programming", 80, PURPLE),
    ("设计能力", "UI/UX / Figma / Creative Direction", 65, TEAL),
    ("工程工具", "Git / Docker / Linux / CLI", 78, PINK),
]

cols = 3
rows = 2
card_w = Inches(3.6)
card_h = Inches(2.2)
gap_x = Inches(0.3)
gap_y = Inches(0.3)
start_x = Inches(0.8)
start_y = Inches(1.5)

for idx, (name, techs, pct, color) in enumerate(skills):
    col = idx % cols
    row = idx // cols
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    # Card background
    add_rect(slide, x, y, card_w, card_h, WHITE, EEEE, Pt(1))

    # Left accent bar
    add_rect(slide, x, y, Pt(4), card_h, color)

    # Skill name
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), card_w - Inches(0.5), Inches(0.4),
                 name, 18, CHARCOAL, True, "Microsoft YaHei")

    # Tech list
    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), card_w - Inches(0.5), Inches(0.5),
                 techs, 13, GRAY_MED, False, "Segoe UI")

    # Progress bar track
    bar_left = x + Inches(0.3)
    bar_top = y + Inches(1.4)
    bar_w = card_w - Inches(0.6)
    add_rect(slide, bar_left, bar_top, bar_w, Pt(6), F0F0F0)

    # Progress bar fill
    fill_w = int(bar_w * pct / 100)
    add_rect(slide, bar_left, bar_top, Emu(fill_w._val), Pt(6), color)

    # Percentage badge
    badge = add_rect(slide, x + card_w - Inches(0.7), y + Inches(1.35), Inches(0.6), Inches(0.35), None, color, Pt(1.5))
    badge.text_frame.paragraphs[0].text = f"{pct}%"
    badge.text_frame.paragraphs[0].font.size = Pt(12)
    badge.text_frame.paragraphs[0].font.color.rgb = color
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.name = "Segoe UI"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.vertical_anchor = MSO_ANCHOR.MIDDLE


# ━━━━━━━━━━━━━━━━━━━━━━━ S04 Project Alpha ━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "项目展示", 28, GRAY_MED, False, "Microsoft YaHei")
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             "Project Alpha — 交互式数据可视化平台", 30, CHARCOAL, True, "Microsoft YaHei")
add_line(slide, Inches(0.8), Inches(1.45), Inches(1.8), Inches(1.45), PURPLE, Pt(3))

# Project thumbnail area
thumb = add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.0), FAFAFA, EEEE, Pt(1))
add_text_box(slide, Inches(1.3), Inches(2.3), Inches(10.5), Inches(1.0),
             "🌍  一个交互式数据可视化平台，支持实时图表渲染和动态数据流处理。",
             20, GRAY_DARK, False, "Microsoft YaHei")

# Tags
tags = [("React", PURPLE), ("D3.js", TEAL), ("WebSocket", PINK)]
for i, (tag, color) in enumerate(tags):
    x = Inches(1.3 + i * 1.5)
    t = add_rect(slide, x, Inches(3.5), Inches(1.3), Inches(0.35), None, color, Pt(1.5))
    t.text_frame.paragraphs[0].text = tag
    t.text_frame.paragraphs[0].font.size = Pt(12)
    t.text_frame.paragraphs[0].font.color.rgb = color
    t.text_frame.paragraphs[0].font.name = "Segoe UI"
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    t.vertical_anchor = MSO_ANCHOR.MIDDLE

# Feature list
features = [
    "支持实时数据流的动态图表渲染",
    "交互式数据探索与多维可视化",
    "WebSocket 双向通信保障低延迟"
]
for i, feat in enumerate(features):
    y = Inches(4.2) + Inches(i * 0.6)
    add_accent_circle(slide, Inches(1.3), y + Inches(0.08), Inches(0.15), PURPLE)
    add_text_box(slide, Inches(1.6), y, Inches(10), Inches(0.5),
                 feat, 18, CHARCOAL, False, "Microsoft YaHei")


# ━━━━━━━━━━━━━━━━━━━━━━━ S05 AI Chat Bot ━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "项目展示", 28, GRAY_MED, False, "Microsoft YaHei")
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             "AI Chat Bot — 智能对话系统", 30, CHARCOAL, True, "Microsoft YaHei")
add_line(slide, Inches(0.8), Inches(1.45), Inches(1.8), Inches(1.45), TEAL, Pt(3))

thumb = add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.0), FAFAFA, EEEE, Pt(1))
add_text_box(slide, Inches(1.3), Inches(2.3), Inches(10.5), Inches(1.0),
             "🤖  基于大语言模型的智能对话系统，支持多轮对话与上下文理解。",
             20, GRAY_DARK, False, "Microsoft YaHei")

tags = [("Python", PURPLE), ("NLP", TEAL), ("FastAPI", PINK)]
for i, (tag, color) in enumerate(tags):
    x = Inches(1.3 + i * 1.5)
    t = add_rect(slide, x, Inches(3.5), Inches(1.3), Inches(0.35), None, color, Pt(1.5))
    t.text_frame.paragraphs[0].text = tag
    t.text_frame.paragraphs[0].font.size = Pt(12)
    t.text_frame.paragraphs[0].font.color.rgb = color
    t.text_frame.paragraphs[0].font.name = "Segoe UI"
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    t.vertical_anchor = MSO_ANCHOR.MIDDLE

features = [
    "基于大语言模型的智能对话引擎",
    "多轮对话上下文理解与记忆",
    "FastAPI 高性能异步接口"
]
for i, feat in enumerate(features):
    y = Inches(4.2) + Inches(i * 0.6)
    add_accent_circle(slide, Inches(1.3), y + Inches(0.08), Inches(0.15), TEAL)
    add_text_box(slide, Inches(1.6), y, Inches(10), Inches(0.5),
                 feat, 18, CHARCOAL, False, "Microsoft YaHei")


# ━━━━━━━━━━━━━━━━━━━━━━━ S06 Dev Toolkit ━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "项目展示", 28, GRAY_MED, False, "Microsoft YaHei")
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
             "Dev Toolkit — 开发者效率工具集", 30, CHARCOAL, True, "Microsoft YaHei")
add_line(slide, Inches(0.8), Inches(1.45), Inches(1.8), Inches(1.45), PINK, Pt(3))

thumb = add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.0), FAFAFA, EEEE, Pt(1))
add_text_box(slide, Inches(1.3), Inches(2.3), Inches(10.5), Inches(1.0),
             "⚙️  开发者效率工具集，包含代码模板生成、API 测试、性能分析等功能。",
             20, GRAY_DARK, False, "Microsoft YaHei")

tags = [("Node.js", PURPLE), ("Electron", TEAL), ("TypeScript", PINK)]
for i, (tag, color) in enumerate(tags):
    x = Inches(1.3 + i * 1.5)
    t = add_rect(slide, x, Inches(3.5), Inches(1.3), Inches(0.35), None, color, Pt(1.5))
    t.text_frame.paragraphs[0].text = tag
    t.text_frame.paragraphs[0].font.size = Pt(12)
    t.text_frame.paragraphs[0].font.color.rgb = color
    t.text_frame.paragraphs[0].font.name = "Segoe UI"
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    t.vertical_anchor = MSO_ANCHOR.MIDDLE

features = [
    "代码模板一键生成",
    "内置 API 测试与调试工具",
    "性能分析与瓶颈检测"
]
for i, feat in enumerate(features):
    y = Inches(4.2) + Inches(i * 0.6)
    add_accent_circle(slide, Inches(1.3), y + Inches(0.08), Inches(0.15), PINK)
    add_text_box(slide, Inches(1.6), y, Inches(10), Inches(0.5),
                 feat, 18, CHARCOAL, False, "Microsoft YaHei")


# ━━━━━━━━━━━━━━━━━━━━━━━ S07 Site Design ━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "网站设计 — 理念与结构", 32, CHARCOAL, True, "Microsoft YaHei")
add_line(slide, Inches(0.8), Inches(1.1), Inches(1.8), Inches(1.1), PURPLE, Pt(3))

# Container
add_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2), FAFAFA, EEEE, Pt(1))

# Left column: design philosophy
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(5), Inches(0.5),
             "设计理念", 22, CHARCOAL, True, "Microsoft YaHei")
add_text_box(slide, Inches(1.2), Inches(2.4), Inches(5), Inches(1.5),
             "相信好的设计加上优雅的代码，\n就能创造出打动人心的产品。",
             18, GRAY_DARK, False, "Microsoft YaHei", PP_ALIGN.LEFT, 1.8)

# Right column: 5 sections
sections = [
    ("Hero — 首页", "粒子网络背景 + 打字机动画", PURPLE),
    ("About — 关于我", "个人介绍与技术理念", TEAL),
    ("Skills — 技能", "六维能力可视化展示", PINK),
    ("Projects — 项目", "作品卡片式展示", PURPLE),
    ("Contact — 联系", "多通道联系方式聚合", TEAL),
]
for i, (sec, desc, color) in enumerate(sections):
    y = Inches(1.8) + Inches(i * 0.95)
    add_accent_circle(slide, Inches(1.3), y + Inches(0.08), Inches(0.15), color)
    add_text_box(slide, Inches(1.6), y, Inches(4.5), Inches(0.35),
                 sec, 16, CHARCOAL, True, "Microsoft YaHei")
    add_text_box(slide, Inches(1.6), y + Inches(0.35), Inches(4.5), Inches(0.3),
                 desc, 13, GRAY_MED, False, "Microsoft YaHei")

# Bottom accent line
add_line(slide, Inches(1.2), Inches(6.4), Inches(11.5), Inches(6.4), PURPLE, Pt(1.5))


# ━━━━━━━━━━━━━━━━━━━━━━━ S08 Tech Highlights ━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "技术实现 — 关键亮点", 32, CHARCOAL, True, "Microsoft YaHei")
add_line(slide, Inches(0.8), Inches(1.1), Inches(1.8), Inches(1.1), PURPLE, Pt(3))

highlights = [
    ("粒子网络背景", "Canvas 绘制 80 个粒子节点，实时连线形成动态网络，鼠标交互产生引力效果", PURPLE),
    ("打字机动画", "多短语循环打字/删除效果，模拟真实打字节奏，营造生动的第一印象", TEAL),
    ("滚动渐入动画", "Intersection Observer 实现元素滚动进入视口时的渐显与技能条动画", PINK),
    ("毛玻璃导航栏", "fixed 定位 + backdrop-filter blur，滚动时自动隐藏/显示导航", PURPLE),
]

cols = 2
rows = 2
card_w = Inches(5.6)
card_h = Inches(2.3)
gap = Inches(0.4)
sx = Inches(0.8)
sy = Inches(1.5)

for idx, (title, desc, color) in enumerate(highlights):
    col = idx % cols
    row = idx // cols
    x = sx + col * (card_w + gap)
    y = sy + row * (card_h + gap)

    add_rect(slide, x, y, card_w, card_h, WHITE, EEEE, Pt(1))

    # Top accent dot
    add_accent_circle(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.2), color)

    add_text_box(slide, x + Inches(0.6), y + Inches(0.2), card_w - Inches(0.8), Inches(0.4),
                 title, 20, CHARCOAL, True, "Microsoft YaHei")
    add_text_box(slide, x + Inches(0.6), y + Inches(0.8), card_w - Inches(0.8), Inches(1.2),
                 desc, 15, GRAY_DARK, False, "Microsoft YaHei", PP_ALIGN.LEFT, 1.7)


# ━━━━━━━━━━━━━━━━━━━━━━━ S09 Future Plans ━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "未来计划 — 持续探索中", 32, CHARCOAL, True, "Microsoft YaHei")
add_line(slide, Inches(0.8), Inches(1.1), Inches(1.8), Inches(1.1), PURPLE, Pt(3))

# Centered container
add_rect(slide, Inches(2.0), Inches(1.8), Inches(9.3), Inches(4.5), FAFAFA, EEEE, Pt(1))

plans = [
    ("深入学习大模型应用开发，探索 AI Agent 方向", PURPLE),
    ("构建更完整的全栈项目，积累工程实践经验", TEAL),
    ("持续打磨前端设计与交互能力", PINK),
    ("保持技术写作与开源贡献", PURPLE),
]

for i, (plan, color) in enumerate(plans):
    y = Inches(2.2) + Inches(i * 0.95)
    add_accent_circle(slide, Inches(2.8), y + Inches(0.08), Inches(0.18), color)
    add_text_box(slide, Inches(3.2), y, Inches(7.8), Inches(0.6),
                 f"→  {plan}", 20, CHARCOAL, False, "Microsoft YaHei")


# ━━━━━━━━━━━━━━━━━━━━━━━ S10 Contact / Thank You ━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)

add_text_box(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(1.0),
             "感谢聆听", 52, CHARCOAL, True, "Microsoft YaHei", PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(2.2), Inches(13.333), Inches(0.6),
             "无论是项目合作、技术交流还是随便聊聊，都欢迎联系我！", 20, GRAY_DARK, False, "Microsoft YaHei", PP_ALIGN.CENTER)

# Contact buttons
contacts = [
    ("Email", "your@email.com", PURPLE),
    ("GitHub", "github.com/yourname", TEAL),
    ("LinkedIn", "linkedin.com/in/yourname", PINK),
    ("Twitter", "@yourname", PURPLE),
]
btn_w = Inches(2.6)
btn_h = Inches(0.6)
gap_btn = Inches(0.4)
total_w = btn_w * 4 + gap_btn * 3
start_x = (W - total_w) / 2

for i, (label, detail, color) in enumerate(contacts):
    x = start_x + i * (btn_w + gap_btn)
    y = Inches(3.3)

    # Button
    btn = add_rect(slide, x, y, btn_w, btn_h, None, color, Pt(2))
    btn.text_frame.word_wrap = True
    btn.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = btn.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

    # Detail below
    add_text_box(slide, x, y + btn_h + Inches(0.1), btn_w, Inches(0.4),
                 detail, 12, GRAY_MED, False, "Segoe UI", PP_ALIGN.CENTER)

# Three-color gradient line at bottom
add_line(slide, Inches(3.5), Inches(5.5), Inches(5.0), Inches(5.5), PURPLE, Pt(3))
add_line(slide, Inches(5.0), Inches(5.5), Inches(6.7), Inches(5.5), TEAL, Pt(3))
add_line(slide, Inches(6.7), Inches(5.5), Inches(8.3), Inches(5.5), PINK, Pt(3))

add_text_box(slide, Inches(0), Inches(6.2), Inches(13.333), Inches(0.5),
             "Q & A", 24, GRAY_MED, True, "Microsoft YaHei", PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━ Save ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
output = os.path.join(os.path.dirname(__file__), "个人作品集分享.pptx")
prs.save(output)
print(f"PPT saved to: {output}")
